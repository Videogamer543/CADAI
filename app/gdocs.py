"""Read Google Docs, Slides, Sheets and Drive files into knowledge-base text.

Why this exists as its own module rather than a branch inside kb_ingest.fetch_doc:
a Google link is not a web page that happens to be hosted at Google. Fetching
https://docs.google.com/presentation/d/<id>/edit gets you an empty HTML shell
and a megabyte of JavaScript -- the slides are drawn in the browser, so the
"text" of that page is about forty characters of chrome. It clears no floor,
carries no headings, and if it ever did index it would answer questions with
menu labels. Every Google document has to be *exported* to something readable
first, and each of the three types exports differently.

No OAuth, no API key, no service account. Google serves the export endpoints to
anonymous requests whenever the document's sharing is set to "Anyone with the
link". That covers essentially every FRC resource anybody actually links to,
and it means this works the moment a link is pasted rather than after a
credentials dance. The cost is that a document shared only with named accounts
cannot be read at all, so the single most important thing in this file is that
such a document produces a sentence naming the setting to change -- see
NotShared. An auth wall does not fail loudly; it returns 200 and an HTML
sign-in page, which is exactly the shape of a successful fetch. Indexing one
puts "Sign in - Google Accounts. Forgot email?" in the corpus under the title
of a document the assistant will then cite.

Export formats, and why each:

  Docs         export?format=html, not txt. webtext.html_to_text turns h1..h6
               into markdown "#" lines, and kb.chunk_text reads those as
               headings and carries them onto every chunk beneath. The txt
               export is the same words with the structure removed, which
               costs the retriever the single best signal it has about what a
               chunk from the middle of a long document is about.

  Slides       export/pptx, read with python-pptx. A deck's structure is its
               meaning: the title of slide 14 is the claim and the body is the
               support, and speaker notes are frequently where the actual
               reasoning was written down. Exporting to txt flattens all of
               that into one column of orphaned phrases. python-pptx is not a
               hard dependency -- see _slides_via_pdf for what happens without
               it.

  Sheets       export?format=xlsx, read with openpyxl, every tab. Rows are
               re-emitted as "Header: value; Header: value" rather than as
               "12 | 0.125 | 6061" because a chunk of bare cells retrieves for
               nothing and reads as noise; see _sheet_text.

  Drive files  uc?export=download. Usually a PDF somebody dropped in a shared
               folder, which app/pdftext.py already reads properly, page
               markers and all.

Slides and Sheets are emitted with "[[page N]]" markers, one per slide or tab,
because kb.chunk_text consumes those into a per-chunk page number and the
assistant cites it. `page_word` travels alongside so slide 7 of a deck is cited
as "(slide 7)" and not "(p. 7)" -- a page number pointing into a document with
no pages is worse than no citation, because the reader goes looking.
"""
from __future__ import annotations

import io
import os
import re
import urllib.parse

from app import pdftext
from app import webtext

UA = webtext.UA
TIMEOUT = 40


class NotShared(ValueError):
    """The document exists but this request is not allowed to read it.

    Its own type, not a bare ValueError, because callers want to say something
    different about it. Every other failure here is "that link did not work";
    this one is "that link works and the fix is two clicks", and a crawl that
    hit forty documents should be able to separate the two at the end.
    """


class NotReadable(ValueError):
    """A Google URL of a kind there is no text to extract from (a Form, a Map)."""


# ---------------------------------------------------------------------------
# URL parsing
# ---------------------------------------------------------------------------
# The /d/<id> and /d/e/<pubid> forms are genuinely different things. /d/<id> is
# the document, and every export endpoint keys off that id. /d/e/<pubid> is a
# *publication* of the document created by File > Publish to the web; it has
# its own opaque id that no export endpoint accepts, and the only readable
# thing at that address is the published HTML itself.
_DOC_RE = re.compile(
    r"https?://docs\.google\.com/(document|presentation|spreadsheets|drawings)"
    r"/(?:u/\d+/)?d/(?P<pub>e/)?(?P<id>[A-Za-z0-9_-]{8,})", re.I)
_DRIVE_RE = re.compile(
    r"https?://drive\.google\.com/(?:file/d/|open\?id=|uc\?[^ ]*[?&]id=)"
    r"(?P<id>[A-Za-z0-9_-]{8,})", re.I)
_DRIVE_FOLDER_RE = re.compile(
    r"https?://drive\.google\.com/drive/(?:u/\d+/)?folders/", re.I)
_FORMS_RE = re.compile(r"https?://docs\.google\.com/forms/", re.I)

_KIND_OF = {"document": "doc", "presentation": "slides",
            "spreadsheets": "sheet", "drawings": "drawing"}


def parse_url(url):
    """{kind, id, published, gid} for a Google URL, or None if it is not one.

    `kind` is one of doc, slides, sheet, drawing, drive, folder, form.
    """
    url = (url or "").strip()
    m = _DOC_RE.match(url)
    if m:
        gid = None
        frag = urllib.parse.urlsplit(url).fragment or ""
        q = urllib.parse.urlsplit(url).query or ""
        gm = re.search(r"gid=(\d+)", frag) or re.search(r"gid=(\d+)", q)
        if gm:
            gid = gm.group(1)
        return {"kind": _KIND_OF[m.group(1).lower()], "id": m.group("id"),
                "published": bool(m.group("pub")), "gid": gid, "url": url}
    if _DRIVE_FOLDER_RE.match(url):
        return {"kind": "folder", "id": "", "published": False, "gid": None,
                "url": url}
    m = _DRIVE_RE.match(url)
    if m:
        return {"kind": "drive", "id": m.group("id"), "published": False,
                "gid": None, "url": url}
    if _FORMS_RE.match(url):
        return {"kind": "form", "id": "", "published": False, "gid": None,
                "url": url}
    return None


def is_google_url(url):
    return parse_url(url) is not None


def share_url(kind, doc_id):
    """The canonical link for a document, used as its identity in the corpus.

    kb.add() replaces on exact URL match, so the same deck reached through an
    /edit link, a /view link and a copy with ?usp=sharing on the end has to
    normalise to one string or a refresh silently triples it.
    """
    if kind == "doc":
        return "https://docs.google.com/document/d/%s/edit" % doc_id
    if kind == "slides":
        return "https://docs.google.com/presentation/d/%s/edit" % doc_id
    if kind == "sheet":
        return "https://docs.google.com/spreadsheets/d/%s/edit" % doc_id
    if kind == "drive":
        return "https://drive.google.com/file/d/%s/view" % doc_id
    return ""


LABEL = {"doc": "Google Doc", "slides": "Google Slides deck",
         "sheet": "Google Sheet", "drive": "Drive file",
         "drawing": "Google Drawing", "folder": "Drive folder",
         "form": "Google Form"}


# ---------------------------------------------------------------------------
# Fetching
# ---------------------------------------------------------------------------
_SIGNIN_RE = re.compile(
    rb"(?i)(accounts\.google\.com/(v\d/)?signin|ServiceLogin|"
    rb"<title>\s*(Sign in|Meld je aan|Request access)|"
    rb"You need (permission|access)|Request access\b)")


def _shared_hint(what, url):
    return (
        "%s is not readable without signing in.\n"
        "        CADAI reads Google documents anonymously, so the document "
        "has to be\n"
        "        link-shared. In the document: Share -> General access -> "
        "\"Anyone with\n"
        "        the link\", role Viewer. Commenter and Editor work too; "
        "\"Restricted\"\n"
        "        does not. Then run this again.\n"
        "        %s" % (what, url))


def _get(url, timeout=TIMEOUT, expect=None):
    """GET with the sign-in wall turned into NotShared instead of into text.

    `expect` is a substring the content-type must contain for the response to
    be what was asked for. It is the sharpest of the three checks: when a
    request for .pptx comes back as text/html, nothing has gone subtly wrong --
    Google has answered a different question, and the answer is a login page.
    """
    import requests
    r = requests.get(url, headers={"User-Agent": UA}, timeout=timeout,
                     allow_redirects=True)
    ctype = (r.headers.get("content-type") or "").lower()
    final = (r.url or "").lower()
    if r.status_code in (401, 403) or "accounts.google.com" in final:
        raise NotShared("sign-in required")
    if r.status_code == 404:
        # Google answers 404 both for a document that does not exist and for
        # one this request may not know exists. They are indistinguishable from
        # here, so say both.
        raise ValueError("not found -- either the link is wrong or the "
                         "document is not shared")
    r.raise_for_status()
    if expect and expect not in ctype:
        if "html" in ctype and _SIGNIN_RE.search(r.content[:8000]):
            raise NotShared("sign-in required")
        raise ValueError("Google returned %s, not the %s export"
                         % (ctype.split(";")[0] or "an unknown type", expect))
    if "html" in ctype and _SIGNIN_RE.search(r.content[:8000]):
        raise NotShared("sign-in required")
    return r


_FNAME_STAR = re.compile(r"filename\*\s*=\s*[\w-]+''([^;]+)", re.I)
_FNAME = re.compile(r'filename\s*=\s*"?([^";]+)', re.I)


def _title_from(resp, fallback=""):
    """The document's real name, out of Content-Disposition.

    Worth the parsing. The export's filename is the title the owner typed, which
    is what somebody scanning a source list will recognise; the alternatives are
    the first line of slide 1 (often "Agenda") and the document id (unreadable).
    """
    cd = resp.headers.get("content-disposition") or ""
    m = _FNAME_STAR.search(cd) or _FNAME.search(cd)
    if not m:
        return fallback
    name = m.group(1).strip()
    try:
        name = urllib.parse.unquote(name)
    except Exception:
        pass
    name = re.sub(r"\.(pptx|xlsx|docx|html?|txt|pdf|csv)$", "", name, flags=re.I)
    return name.strip() or fallback


# ---------------------------------------------------------------------------
# Google Docs
# ---------------------------------------------------------------------------
def read_doc(doc_id, timeout=TIMEOUT):
    url = "https://docs.google.com/document/d/%s/export?format=html" % doc_id
    r = _get(url, timeout=timeout)
    r.encoding = r.encoding or "utf-8"
    title, text = webtext.html_to_text(r.text)
    return {"title": _title_from(r, title), "text": text,
            "pages": 0, "page_word": ""}


def read_published(url, timeout=TIMEOUT):
    """A File > Publish to the web address, which has no export endpoint.

    Only Docs publish to real HTML. A published deck is the same JavaScript
    canvas as the editor and a published sheet is a table rendered by script,
    so both come back nearly empty -- and an almost-empty document is the one
    failure worth naming precisely, because "added 0 documents" sends people
    looking at their sharing settings when the settings are fine.
    """
    r = _get(url, timeout=timeout)
    r.encoding = r.encoding or "utf-8"
    title, text = webtext.html_to_text(r.text)
    if len((text or "").strip()) < 200:
        raise ValueError(
            "a published-to-web link, and almost no text came back.\n"
            "        Published Slides and Sheets are drawn by JavaScript, so "
            "there is\n"
            "        nothing to read at that address. Open the document and "
            "use the\n"
            "        normal Share link instead -- the one with /d/<id>/edit in "
            "it.")
    return {"title": _title_from(r, title), "text": text,
            "pages": 0, "page_word": ""}


# ---------------------------------------------------------------------------
# Google Slides
# ---------------------------------------------------------------------------
def _shape_text(shape, out):
    """Append one shape's text, descending into groups and tables."""
    try:
        if shape.shape_type is not None and str(shape.shape_type).startswith("GROUP"):
            for sub in sorted(shape.shapes,
                              key=lambda s: ((s.top or 0), (s.left or 0))):
                _shape_text(sub, out)
            return
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                cells = [" ".join((c.text or "").split()) for c in row.cells]
                cells = [c for c in cells if c]
                if cells:
                    out.append(" | ".join(cells))
            return
    except Exception:
        pass
    try:
        if not shape.has_text_frame:
            return
    except Exception:
        return
    for para in shape.text_frame.paragraphs:
        line = " ".join("".join(r.text or "" for r in para.runs).split())
        if not line:
            continue
        # A bullet in a deck is a bullet in the text. kb.chunk_text splits on
        # blank lines, so leaving the lines of one text box unseparated keeps
        # a slide's body together as one block instead of scattering it.
        if para.level:
            line = "  " * para.level + "- " + line
        out.append(line)


def _slides_via_pptx(blob, timeout=None):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob))
    parts, n = [], 0
    for i, slide in enumerate(prs.slides, 1):
        title, title_id = "", None
        try:
            if slide.shapes.title is not None:
                title = " ".join((slide.shapes.title.text or "").split())
                # By shape_id, not by identity. python-pptx builds a fresh
                # proxy object on every attribute access, so `shape is
                # slide.shapes.title` is False for the title shape itself --
                # which silently prints every slide's title twice, once as the
                # heading and once as the first line of its own body.
                title_id = slide.shapes.title.shape_id
        except Exception:
            pass
        body = []
        # Reading order, not XML order. A deck's shapes are stored in the order
        # they were drawn, which after any amount of editing is arbitrary; top
        # to bottom then left to right is what a person reading the slide does.
        try:
            shapes = sorted(slide.shapes,
                            key=lambda s: ((s.top or 0), (s.left or 0)))
        except Exception:
            shapes = list(slide.shapes)
        for shape in shapes:
            try:
                if title_id is not None and shape.shape_id == title_id:
                    continue
            except Exception:
                pass
            _shape_text(shape, body)
        # Slide numbers and footers export as their own little text boxes and
        # would otherwise appear as a paragraph reading "14".
        body = [b for b in body if not re.fullmatch(r"[\d\s/|.-]{0,6}", b)]
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = " ".join(
                    (slide.notes_slide.notes_text_frame.text or "").split())
        except Exception:
            pass
        if not (title or body or notes):
            continue
        n = i
        parts.append(pdftext.PAGE_MARK % i)
        # Always a heading, even when the slide has no title of its own.
        # chunk_text carries the nearest heading above onto every chunk, so a
        # titleless slide would otherwise be filed under the previous slide's
        # subject and retrieve as if it were about that.
        parts.append("## " + (title or ("Slide %d" % i)))
        if body:
            parts.append("\n".join(body))
        if notes:
            # Labelled, because notes are the presenter's reasoning and body
            # text is the claim; merged silently they read as one voice.
            parts.append("Speaker notes: " + notes)
    return "\n\n".join(parts), n


def _slides_via_pdf(blob):
    """Fallback when python-pptx is missing: the PDF export, read by pdftext.

    Costs the speaker notes and the title/body distinction, keeps the words and
    the per-slide page numbers. Worth having as a fallback rather than a hard
    failure -- python-pptx pulls lxml, and a machine where that wheel does not
    build should still be able to ingest a deck.
    """
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as fh:
        fh.write(blob)
        tmp = fh.name
    try:
        info = pdftext.read(tmp)
        return info["text"], info.get("n_pages") or len(info.get("pages") or [])
    finally:
        try:
            os.unlink(tmp)
        except OSError:
            pass


def read_slides(doc_id, timeout=TIMEOUT):
    try:
        import pptx  # noqa: F401
        have_pptx = True
    except ImportError:
        have_pptx = False
    if have_pptx:
        r = _get("https://docs.google.com/presentation/d/%s/export/pptx"
                 % doc_id, timeout=timeout)
        text, n = _slides_via_pptx(r.content)
        note = ""
    else:
        r = _get("https://docs.google.com/presentation/d/%s/export/pdf"
                 % doc_id, timeout=timeout)
        text, n = _slides_via_pdf(r.content)
        note = ("read from the PDF export -- speaker notes were not included. "
                "For those: pip install python-pptx")
    return {"title": _title_from(r, ""), "text": text, "pages": n,
            "page_word": "slide", "note": note}


# ---------------------------------------------------------------------------
# Google Sheets
# ---------------------------------------------------------------------------
MAX_ROWS = 3000
MAX_COLS = 40


def _cell(v):
    if v is None:
        return ""
    if isinstance(v, float) and v == int(v):
        v = int(v)
    return " ".join(str(v).split())


def _sheet_text(rows, name):
    """One tab's rows as retrievable prose.

    The decision that matters: when the tab looks tabular, each row is emitted
    as "Part: P4019; Material: 6061-T6; Thickness: 0.125" rather than as
    "P4019 | 6061-T6 | 0.125". A chunk of the second kind is 900 characters of
    values with their meanings 80 rows above, outside the window -- it matches
    no question anybody asks and, if it is ever retrieved, the assistant cannot
    tell which column is which either. Repeating the header on every row costs
    maybe three times the characters and is the whole difference between a
    spreadsheet being in the corpus and being usable from it.

    "Looks tabular" is deliberately conservative: two or more named columns and
    at least three data rows. Below that it is a scratch pad or a form, and
    pasting a header onto it invents structure that is not there.
    """
    rows = [[_cell(c) for c in row[:MAX_COLS]] for row in rows]
    rows = [r for r in rows if any(r)]
    if not rows:
        return ""
    head = rows[0]
    named = [h for h in head if h]
    tabular = len(named) >= 2 and len(rows) >= 4
    out = ["## " + name]
    if tabular:
        out.append("Columns: " + ", ".join(named))
        for row in rows[1:]:
            bits = []
            for i, val in enumerate(row):
                if not val:
                    continue
                key = head[i] if i < len(head) and head[i] else ""
                bits.append(("%s: %s" % (key, val)) if key else val)
            if bits:
                out.append("; ".join(bits))
    else:
        for row in rows:
            vals = [v for v in row if v]
            if vals:
                out.append(" | ".join(vals))
    return "\n\n".join(out)


def _sheets_via_xlsx(blob):
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(blob), data_only=True, read_only=True)
    parts, n, clipped = [], 0, []
    for i, ws in enumerate(wb.worksheets, 1):
        rows, over = [], False
        for j, row in enumerate(ws.iter_rows(values_only=True)):
            if j >= MAX_ROWS:
                over = True
                break
            rows.append(list(row))
        body = _sheet_text(rows, ws.title or ("Sheet %d" % i))
        # A workbook usually carries one or two empty tabs somebody made and
        # never filled. Emitting them costs a page number and a heading that
        # retrieves for the tab's name and then has nothing under it.
        if len(body.strip().split("\n", 1)) < 2:
            continue
        n = i
        parts.append(pdftext.PAGE_MARK % i)
        parts.append(body)
        if over:
            clipped.append(ws.title or ("sheet %d" % i))
    try:
        wb.close()
    except Exception:
        pass
    note = ("first %d rows only of: %s" % (MAX_ROWS, ", ".join(clipped))
            if clipped else "")
    return "\n\n".join(parts), n, note


def _sheets_via_csv(doc_id, gid, timeout):
    """Fallback without openpyxl: CSV, which is one tab and only one tab.

    Said out loud by the caller rather than left to be discovered, because a
    workbook whose second tab holds the actual data indexes as a success here.
    """
    import csv
    u = ("https://docs.google.com/spreadsheets/d/%s/export?format=csv" % doc_id)
    if gid:
        u += "&gid=%s" % gid
    r = _get(u, timeout=timeout)
    r.encoding = r.encoding or "utf-8"
    rows = list(csv.reader(io.StringIO(r.text)))[:MAX_ROWS]
    body = _sheet_text(rows, "Sheet")
    return r, ("%s\n\n%s" % (pdftext.PAGE_MARK % 1, body) if body else ""), 1


def read_sheet(doc_id, gid=None, timeout=TIMEOUT):
    try:
        import openpyxl  # noqa: F401
        have = True
    except ImportError:
        have = False
    if have:
        r = _get("https://docs.google.com/spreadsheets/d/%s/export?format=xlsx"
                 % doc_id, timeout=timeout)
        text, n, note = _sheets_via_xlsx(r.content)
    else:
        r, text, n = _sheets_via_csv(doc_id, gid, timeout)
        note = ("read as CSV -- only the first tab. For every tab: "
                "pip install openpyxl")
    return {"title": _title_from(r, ""), "text": text, "pages": n,
            "page_word": "tab", "note": note}


# ---------------------------------------------------------------------------
# Drive files
# ---------------------------------------------------------------------------
def read_drive_file(file_id, timeout=TIMEOUT):
    """Whatever somebody parked in a shared Drive folder. Usually a PDF."""
    import tempfile
    r = _get("https://drive.google.com/uc?export=download&id=%s" % file_id,
             timeout=timeout)
    blob = r.content
    ctype = (r.headers.get("content-type") or "").lower()
    if blob[:4] == b"%PDF" or "pdf" in ctype:
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as fh:
            fh.write(blob)
            tmp = fh.name
        try:
            info = pdftext.read(tmp)
            return {"title": _title_from(r, info.get("title") or ""),
                    "text": info["text"],
                    "pages": info.get("n_pages") or 0, "page_word": "p.",
                    "info": info}
        finally:
            try:
                os.unlink(tmp)
            except OSError:
                pass
    if blob[:4] == b"PK\x03\x04":
        name = _title_from(r, "")
        if name and re.search(r"pptx?$", ctype) or "presentation" in ctype:
            text, n = _slides_via_pptx(blob)
            return {"title": name, "text": text, "pages": n,
                    "page_word": "slide"}
        if "sheet" in ctype or "excel" in ctype:
            text, n, note = _sheets_via_xlsx(blob)
            return {"title": name, "text": text, "pages": n,
                    "page_word": "tab", "note": note}
    if "html" in ctype or "text/" in ctype:
        r.encoding = r.encoding or "utf-8"
        title, text = webtext.html_to_text(r.text)
        return {"title": _title_from(r, title), "text": text,
                "pages": 0, "page_word": ""}
    raise ValueError("a Drive file of a type this cannot read (%s)"
                     % (ctype.split(";")[0] or "unknown"))


# ---------------------------------------------------------------------------
# The one entry point
# ---------------------------------------------------------------------------
def fetch(url, timeout=TIMEOUT):
    """Read any Google URL. Returns the shape kb_ingest.fetch_doc returns.

    {title, text, info, is_pdf, pages, page_word, note, gkind}
    """
    g = parse_url(url)
    if not g:
        raise ValueError("not a Google Docs/Slides/Sheets/Drive URL")
    kind, gid = g["kind"], g["id"]

    if kind == "folder":
        raise NotReadable(
            "a Drive folder, not a document. Open it and paste the links to "
            "the\n        files inside -- several links at once is fine.")
    if kind == "form":
        raise NotReadable("a Google Form. There is nothing in it to read.")
    if kind == "drawing":
        raise NotReadable("a Google Drawing -- a picture, with no text layer.")

    try:
        if g["published"]:
            out = read_published(url, timeout=timeout)
            out["gkind"] = kind
        elif kind == "doc":
            out = read_doc(gid, timeout=timeout)
            out["gkind"] = "doc"
        elif kind == "slides":
            out = read_slides(gid, timeout=timeout)
            out["gkind"] = "slides"
        elif kind == "sheet":
            out = read_sheet(gid, g.get("gid"), timeout=timeout)
            out["gkind"] = "sheet"
        else:
            out = read_drive_file(gid, timeout=timeout)
            out["gkind"] = "drive"
    except NotShared:
        raise NotShared(_shared_hint("this %s" % LABEL.get(kind, "document"),
                                     share_url(kind, gid) or url))
    out.setdefault("info", None)
    out.setdefault("note", "")
    out.setdefault("pages", 0)
    out.setdefault("page_word", "")
    out["is_pdf"] = bool(out.get("info"))
    if not out.get("title"):
        out["title"] = "%s %s" % (LABEL.get(kind, "Google document"), gid[:8])
    out["url"] = share_url(kind, gid) or url
    return out
